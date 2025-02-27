from pptx import Presentation
fo = open("DLcontent.txt", "w",encoding='utf8', errors='ignore')

# for x in range(9):
prs = Presentation('Digital Logic.pptx')
# text_runs will be populated with a list of strings,
# one for each text run in presentation
for slide in prs.slides:
	list_of_elements = []
	for shape in slide.shapes:
		if not shape.has_text_frame:
			continue
		for paragraph in shape.text_frame.paragraphs:
			line = ''
			for run in paragraph.runs:
				line += run.text	
			list_of_elements.append(line)
		
	for elements in list_of_elements:
		print(elements)
		fo.write(elements +'\n')
	fo.write('\n')
